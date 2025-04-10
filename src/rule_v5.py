# ruleEngine.py
import json
import csv
from tkinter import Image
import matplotlib.pyplot as plt
import logging
from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import SimpleDocTemplate, Image as ReportLabImage, Spacer, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
import fitz  # PyMuPDF
import os
from PIL import Image as PILImage
from reportlab.pdfgen import canvas
from difflib import SequenceMatcher
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph
from reportlab.lib.enums import TA_CENTER
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Load OCR Results
ocr_file = "output/ocr_results.json"
with open(ocr_file, "r") as f:
    ocr_data = json.load(f)
    
sample_no=1
THRESHOLD = 0.3  # Adjust as needed
HEADING_SIMILARITY_THRESHOLD = 0.8 # Adjust for heading similarity
MIN_LCS_LENGTH = 6 # Minimum length of longest common substring for considering same heading
PDF_PATH = f"input/sample{sample_no}.pdf"
input_csv = f"input/sample{sample_no}.csv"
control_csv = f"input/control_sample{sample_no}.csv"
OUTPUT_COMBINED_DIR = f"output/combined_error_pages_sample{sample_no}"
OUTPUT_PPTX = "output/error_pages.pptx"
csv_file = f"output/processed_sample{sample_no}.csv"
os.makedirs(OUTPUT_COMBINED_DIR, exist_ok=True)


def is_similar(str1, str2):
    """Check if two strings are similar based on SequenceMatcher."""
    if not str1 or not str2:
        return False
    return SequenceMatcher(None, str1.lower(), str2.lower()).ratio() >= HEADING_SIMILARITY_THRESHOLD

def longest_common_substring_length(s1, s2):
    """Calculate the length of the longest common substring."""
    m = len(s1)
    n = len(s2)
    dp = [[0] * (n + 1) for _ in range(m + 1)]
    max_len = 0
    for i in range(1, m + 1):
        for j in range(1, n + 1):
            if s1[i - 1].lower() == s2[j - 1].lower():
                dp[i][j] = dp[i - 1][j - 1] + 1
                max_len = max(max_len, dp[i][j])
            else:
                dp[i][j] = 0
    return max_len

def group_pages(ocr_data):
    """Groups pages with priority on 'page_info_strict', then date continuity, then headings, then other continuity."""
    grouped_sections = []
    current_group = []
    n_pages = len(ocr_data["pages"])
    i = 0

    while i < n_pages:
        page_info = ocr_data["pages"][i]
        page_number = page_info["page_number"]
        heading = page_info["signals"]["heading"]
        dates = page_info["signals"]["date"] if page_info["signals"]["date"] else []
        page_info_strict = page_info["signals"]["page_info_strict"]
        continuity_weight = page_info["continuity_weight"]
        continuity_reasons = page_info["continuity_reasons"]

        # First Priority: "Page 1 of X" with X >= 2
        if page_info_strict and page_info_strict[0] == 1 and page_info_strict[1] >= 2:
            if current_group and len(current_group) > 1:
                grouped_sections.append(current_group)
            current_group = []
            num_pages_in_section = page_info_strict[1]
            end_page = min(page_number + num_pages_in_section, n_pages + 1)
            group = [(ocr_data["pages"][j]["page_number"],
                      ocr_data["pages"][j]["signals"]["heading"],
                      ocr_data["pages"][j]["signals"]["date"])
                     for j in range(i, end_page)]
            grouped_sections.append(group)
            i = end_page
            continue

        # Second Priority: Consecutive pages with at least one common date
        if dates and current_group and current_group[-1][2]:
            last_group_dates = set(current_group[-1][2])
            current_page_dates = set(dates)
            if any(date in last_group_dates for date in current_page_dates):
                current_group.append((page_number, heading, dates))
            else:
                if current_group:
                    grouped_sections.append(current_group)
                current_group = [(page_number, heading, dates)]
        elif dates and not current_group:
            current_group = [(page_number, heading, dates)]

        # Third Priority: Adjacent pages with the same heading (LCS length > MIN_LCS_LENGTH)
        elif heading and i + 1 < n_pages and ocr_data["pages"][i + 1]["signals"]["heading"]:
            next_heading = ocr_data["pages"][i + 1]["signals"]["heading"]
            if longest_common_substring_length(heading, next_heading) > MIN_LCS_LENGTH:
                if current_group and len(current_group) > 1:
                    grouped_sections.append(current_group)
                grouped_sections.append([(page_number, heading, dates)])
                current_group = []
                i += 1 # Skip the next page as it's now in its own group
                continue

        # Fourth Priority: Continuity Weight
        elif continuity_weight > THRESHOLD and continuity_reasons:
            current_group.append((page_number, heading, dates))

        # Fifth Priority: Heading Similarity
        elif heading and i + 1 < n_pages and is_similar(heading, ocr_data["pages"][i + 1]["signals"]["heading"]):
            if current_group and len(current_group) > 1:
                grouped_sections.append(current_group)
            grouped_sections.append([(page_number, heading, dates)])
            current_group = []
            i += 1 # Skip the next page
            continue
        else:
            # Lowest Priority: Individual pages
            if current_group:
                grouped_sections.append(current_group)
            current_group = [(page_number, heading, dates)]

        i += 1

    if current_group:
        grouped_sections.append(current_group)

    final_sections = []
    for group in grouped_sections:
        if group:
            final_sections.append(group)

    return final_sections

def save_to_csv(grouped_sections):
    """Save grouped sections into a CSV."""
    
    with open(csv_file, "w", newline="") as f:
        writer = csv.writer(f)
        for group in grouped_sections:
            if group:
                start_page = group[0][0]
                end_page = group[-1][0]
                heading = " | ".join(set(entry[1] for entry in group if entry[1] is not None))
                start_dates = [date for entry in group for date in (entry[2] if entry[2] else [])]
                end_dates = [date for entry in reversed(group) for date in (entry[2] if entry[2] else [])]
                start_date = start_dates[0] if start_dates else "No Date"
                end_date = end_dates[0] if end_dates else "No Date"
                writer.writerow([start_page, end_page, heading, start_date, end_date])

    print(f"✅ Processed CSV saved at {csv_file}")
    return csv_file

def load_input_csv(file_path, y_offset=0):
    """Load page groups from a CSV."""
    page_values = {}
    step_data = [(0, y_offset)]
    with open(file_path, "r") as f:
        reader = csv.reader(f)
        next(reader, None)
        for row in reader:
            if len(row) >= 2:
                try:
                    start_page = int(row[0])
                    end_page = int(row[1])
                    for page in range(start_page, end_page + 1):
                        page_values[page] = 1
                    step_data.extend([(start_page, y_offset), (start_page, y_offset + 1), (end_page, y_offset + 1),
                                      (end_page, y_offset)])
                except ValueError as e:
                    logging.warning(f"Skipping invalid row in CSV: {row} - {e}")
    return step_data, page_values


def extract_ranges(file_path,controlvector=None):
    """Extract page ranges (start, end) from a CSV file."""
    ranges = []
    with open(file_path, "r") as f:
        reader = csv.reader(f)
        next(reader, None)  # Skip header if present
        for row in reader:
            if len(row) >= 2:
                try:
                    start_page = int(row[0])
                    if controlvector:
                        controlvector[start_page-1] = 1
                    # Check if the end page is a valid integer
                    end_page = int(row[1])
                    ranges.append((start_page, end_page))
                except ValueError as e:
                    logging.warning(f"Invalid row in CSV: {row} - {e}")
    return ranges,controlvector


def is_page_in_ranges(page, ranges):
    """Check if a given page number falls within any of the provided ranges."""
    for start, end in ranges:
        if start <= page <= end:
            return True
    return False


def calculate_error_percentage(processed_file, input_file, control_file,controlvector,processedvector):
    """Calculate error percentage by comparing page ranges."""

    processed_ranges,processedvector = extract_ranges(processed_file,processedvector)
    input_ranges,dummy = extract_ranges(input_file,None)
    control_ranges,controlvector = extract_ranges(control_file,controlvector)

    all_relevant_pages = set()
    for start, end in processed_ranges + input_ranges + control_ranges:
        all_relevant_pages.update(range(start, end + 1))

    if not all_relevant_pages:
        return 0, 0, set(), set(), set(), set()

    error_input_pages = set()
    for page in all_relevant_pages:
        in_processed = is_page_in_ranges(page, processed_ranges)
        in_input = is_page_in_ranges(page, input_ranges)
        if in_processed != in_input:
            error_input_pages.add(page)

    error_control_pages = set()
    for page in all_relevant_pages:
        in_processed = is_page_in_ranges(page, processed_ranges)
        in_control = is_page_in_ranges(page, control_ranges)
        if in_processed != in_control:
            error_control_pages.add(page)

    total_relevant = len(all_relevant_pages)
    percentage_input = (len(error_input_pages) / total_relevant) * 100 if total_relevant else 0
    percentage_control = (len(error_control_pages) / total_relevant) * 100 if total_relevant else 0

    print(f"\n✅ Error Percentage (Input): {percentage_input:.2f}%")
    print(f"\n✅ Error Percentage (Control): {percentage_control:.2f}%\n")
    return percentage_input, percentage_control, control_ranges, processed_ranges, input_ranges, control_ranges, error_input_pages, error_control_pages,controlvector,processedvector


def plot_comparison_chart(processed_file, input_file, control_file, error_input_pages=set(), error_control_pages=set()):
    """Compare processed grouping with input and control."""
    processed_data, _ = load_input_csv(processed_file, y_offset=0)
    input_data, _ = load_input_csv(input_file, y_offset=4)
    control_data, _ = load_input_csv(control_file, y_offset=8)

    plt.figure(figsize=(12, 8))
    plt.plot(*zip(*processed_data), marker='o', linestyle='-', color='b', label="Processed")
    plt.plot(*zip(*input_data), marker='o', linestyle='-', color='gray', label="Input")
    plt.plot(*zip(*control_data), marker='o', linestyle='-', color='g', label="Control")

    if error_input_pages:
        plt.scatter(list(error_input_pages), [4.5] * len(error_input_pages), color='r', marker='x', s=100, label="Input Errors")
    if error_control_pages:
        plt.scatter(list(error_control_pages), [8.5] * len(error_control_pages), color='orange', marker='x', s=100,
                    label="Control Errors")

    plt.xlabel("Page Numbers")
    plt.ylabel("Grouping")
    plt.title("Comparison of Page Grouping")
    plt.yticks([0.5, 4.5, 8.5], ["Processed", "Input", "Control"])
    plt.grid(True)
    plt.legend()
    plt.savefig("output/group_comparison_chart.png")
    plt.show()
    print("📊 Graph saved at output/group_comparison_chart.png")
    return set(), set() # Return empty sets as error calculation is now separate


def get_page_image(doc, page_num, img_path):
    """Extracts and saves a page as an image."""
    try:
        if 1 <= page_num <= len(doc):
            page = doc[page_num - 1]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Increase resolution
            pix.save(img_path)
            return True
        return False
    except Exception as e:
        logging.error(f"Error extracting image for page {page_num}: {e}")
        return False



def create_error_pdf(error_pages,error_string, pdf_path, output_pdf_path):
    """Creates a PPTX presentation containing each error page with context (wide image)."""
    prs = Presentation()
    prs.slide_width = Inches(16)  # Landscape
    prs.slide_height = Inches(9)

    doc = fitz.open(pdf_path)
    num_pages = len(doc)
    padding = 10
    generated_combined_images = []

    for x, error_page_num in enumerate(error_pages):
        logging.info(f"Processing error page for PPTX: {error_page_num}")

        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        pil_images = []
        widths = []
        max_height = 0
        temp_files_to_remove = []
        combined_image_path = None

        for i in range(max(1, error_page_num - 1), min(num_pages + 1, error_page_num + 2)):
            temp_path = os.path.join(OUTPUT_COMBINED_DIR, f"temp_page_{error_page_num}_{i}.png")
            if get_page_image(doc, i, temp_path):
                try:
                    img = PILImage.open(temp_path)
                    pil_images.append(img)
                    widths.append(img.width)
                    max_height = max(max_height, img.height)
                    temp_files_to_remove.append(temp_path)
                except FileNotFoundError:
                    logging.error(f"Could not open image: {temp_path}")
            else:
                logging.warning(f"Could not extract image for page {i} (near error page {error_page_num}).")
                blank_img = PILImage.new('RGB', (100, 100), color='lightgray')
                pil_images.append(blank_img)
                widths.append(blank_img.width)
                max_height = max(max_height, blank_img.height)

        if pil_images:
            total_width = sum(widths) + padding * (len(pil_images) - 1)
            combined_image = PILImage.new('RGB', (total_width, max_height), color='white')
            x_offset = 0
            for img in pil_images:
                combined_image.paste(img, (x_offset, 0))
                x_offset += img.width + padding

            combined_image_path = os.path.join(OUTPUT_COMBINED_DIR,
                                                f"{error_string[x]}___combined_error_page_{error_page_num}.png")
            try:
                combined_image.save(combined_image_path)
                logging.info(f"✅ Combined image saved for PPTX page {error_page_num} at {combined_image_path}")
                generated_combined_images.append(combined_image_path)

                # Add heading (error page number) at the top
                left = Inches(0.5)
                top = Inches(0.5)
                width = prs.slide_width - Inches(1)
                height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = txBox.text_frame
                p = text_frame.add_paragraph()
                p.text = f"Error Page: {error_page_num}"
                p.font.size = Inches(0.5)
                p.font.bold = True
                text_frame.vertical_anchor = MSO_ANCHOR.TOP  # Align heading to top
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

                # Add combined image below the heading, stretching wide
                image = PILImage.open(combined_image_path)
                img_width_px, img_height_px = image.size
                image_aspect_ratio = img_width_px / img_height_px

                available_width = prs.slide_width - Inches(1)
                image_width_on_slide = available_width
                image_height_on_slide = image_width_on_slide / image_aspect_ratio

                # Position the image below the heading
                image_top = Inches(1.5) # Adjust as needed based on heading height

                if image_height_on_slide > prs.slide_height - Inches(2):
                    image_height_on_slide = prs.slide_height - Inches(2)
                    image_width_on_slide = image_height_on_slide * image_aspect_ratio
                    image_left = (prs.slide_width - image_width_on_slide) / 2
                    slide.shapes.add_picture(combined_image_path, image_left, image_top, width=image_width_on_slide, height=image_height_on_slide)
                else:
                    image_left = Inches(0.5)
                    slide.shapes.add_picture(combined_image_path, image_left, image_top, width=image_width_on_slide, height=image_height_on_slide)


            except Exception as e:
                logging.error(f"Error adding combined image to PPTX for page {error_page_num}: {e}")

            for img in pil_images:
                img.close()

        for temp_file in temp_files_to_remove:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    prs.save(OUTPUT_PPTX)
    doc.close()
    print(f"✅ Error pages combined into PPTX (wide image): {OUTPUT_PPTX}")


if __name__ == "__main__":
    ocr_results = None
    n_pages = len(ocr_data["pages"])
    controlvector= [0]*n_pages
    processedvector= [0]*n_pages
    errorCounter=0
    errorlist=[]
    errorString=[]
    
    if os.path.exists(ocr_file):
        with open(ocr_file, "r") as f:
            ocr_results = json.load(f)
    else:
        doc = fitz.open(PDF_PATH)
        images = [Image.frombytes("RGB", [page.get_pixmap().width, page.get_pixmap().height],
                                    page.get_pixmap().samples) for page in doc]
        from signal_v4 import perform_ocr, save_results
        ocr_results = perform_ocr(images)
        save_results(ocr_results)
        doc.close()
        print(f"OCR processing completed and saved to {ocr_file}")

    if ocr_results:
        grouped_sections = group_pages(ocr_results)
        csv_file = save_to_csv(grouped_sections)

        error_percentage_input, error_percentage_control, control_ranges, processed_ranges, input_ranges, control_ranges, error_input_pages, error_control_pages,controlvector,processedvector = calculate_error_percentage(
            csv_file, input_csv, control_csv,controlvector,processedvector)

        #print("\n🚨 Pages with Errors (Input):", sorted(list(error_input_pages)))
        #print("\n🚨 Pages with Errors (Control):", sorted(list(error_control_pages)))
        for i in range(0, n_pages):
                print(f"Control Vector : Page {i+1}: {controlvector[i]}\t\tProcessed Vector : Page {i+1}: {processedvector[i]}\t\t Error = {controlvector[i]^processedvector[i]}")
                if controlvector[i]^processedvector[i] == 1:
                    errorCounter+=1
                    errorString.append("DC") if controlvector[i] else  errorString.append("CD")
                    print(f"Error on Page {i+1}")
                    errorlist.append(i+1)
        print(f"Error Pages : {errorlist}")
        print(f"Total Error Percentage : {errorCounter/n_pages*100:.2f}%")
        #plot_comparison_chart(csv_file, input_csv, control_csv, error_input_pages, error_control_pages)
        print(f"✅ Combined error page images saved in: {OUTPUT_COMBINED_DIR}")
        create_error_pdf(errorlist,errorString, PDF_PATH, OUTPUT_PPTX)
        print(f"✅ Combined error pages PDF saved at {OUTPUT_PPTX}")
        # create_error_pdf(error_input_pages, PDF_PATH, OUTPUT_ERROR_PDF.replace(".pdf", "_input_errors.pdf"))